# src/spindle/processors/pptx_setter_processor.py

from spindle.abstracts import AbstractProcessor
from typing import Any, Dict
from pptx import Presentation
from pptx.util import Inches
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape

__all__ = ['PPTXSetterProcessor']


class PPTXSetterProcessor(AbstractProcessor):
    def _preprocess(self, content: Presentation, **kwargs) -> Presentation:
        return content

    def _extract_content(self, presentation: Presentation, **kwargs) -> Dict[str, Any]:
        return {'presentation': presentation}

    def _main_process(self, data: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        presentation = data['presentation']
        text = kwargs.get('text')
        is_notes = kwargs.get('is_notes', False)
        slide_index = kwargs.get('slide_index')

        if slide_index is not None:
            self._set_single_slide(presentation, slide_index, text, is_notes)
        else:
            self._set_all_slides(presentation, text, is_notes)

        return {'presentation': presentation}

    def _postprocess(self, data: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        return data

    def _set_single_slide(self, presentation: Presentation, index: int, content: str, is_notes: bool):
        if 0 <= index < len(presentation.slides):
            slide = presentation.slides[index]
            if is_notes:
                self._set_slide_notes(slide, content)
            else:
                self._set_slide_content(slide, content)
        else:
            raise ValueError(f"Slide index {index} is out of range. Presentation has {len(presentation.slides)} slides.")

    def _set_all_slides(self, presentation: Presentation, content: str, is_notes: bool):
        for slide in presentation.slides:
            if is_notes:
                self._set_slide_notes(slide, content)
            else:
                self._set_slide_content(slide, content)

    def _set_slide_content(self, slide, content: str):
        if content:
            text_frame = self._get_or_create_text_frame(slide)
            text_frame.text = content

    def _set_slide_notes(self, slide, notes: str):
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes

    def _get_or_create_text_frame(self, slide) -> Shape:
        for shape in slide.shapes:
            if isinstance(shape, BaseShape) and hasattr(shape, 'text_frame'):
                return shape.text_frame

        # If no suitable shape is found, create a new text box
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        return txBox.text_frame
