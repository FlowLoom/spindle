import logging
from spindle.abstracts import AbstractProcessor
from typing import Any, Dict, List, Optional, Union
from pptx import Presentation

__all__ = ['PPTXProcessor']

logger = logging.getLogger(__name__)

class PPTXProcessor(AbstractProcessor):
    def _preprocess(self, content: Presentation, **kwargs) -> Presentation:
        logger.info(f"Processing presentation with {len(content.slides)} slides")
        return content

    def _extract_content(self, presentation: Presentation, **kwargs) -> Dict[str, Any]:
        only_metadata = kwargs.get('only_metadata', False)
        if only_metadata:
            return {'metadata': self._extract_metadata(presentation)}

        slide_range = kwargs.get('slide_range')
        metadata = kwargs.get('no_metadata')
        if slide_range is not None:
            return self._extract_slide_range(presentation, slide_range, metadata)
        else:
            return self._extract_all_slides(presentation, **kwargs)

    def _main_process(self, data: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        return data

    def _postprocess(self, data: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        if kwargs.get('only_metadata'):
            return data

        if 'slides' in data:
            if kwargs.get('only_content'):
                data['slides'] = [{'slide_number': slide['slide_number'], 'content': slide['content']} for slide in data['slides']]
            elif kwargs.get('only_notes'):
                data['slides'] = [{'slide_number': slide['slide_number'], 'notes': slide['notes']} for slide in data['slides']]

        if kwargs.get('no_metadata'):
            data.pop('metadata', None)

        return data

    def _extract_slide_range(self, presentation: Presentation, slide_range: Union[int, tuple], metadata: Optional[Union[bool, None]]) -> Dict[str, Any]:

        if isinstance(slide_range, int):
            slide_range = (slide_range, slide_range + 1)

        start, end = slide_range
        if start < 0 or end > len(presentation.slides) or start >= end:
            raise ValueError(f"Invalid slide range. Presentation has {len(presentation.slides)} slides.")

        slides = []
        for i in range(start, end):
            slide = presentation.slides[i]
            slides.append({
                'slide_number': i + 1,
                'content': self._extract_slide_content(slide),
                'notes': self._extract_speaker_notes(slide)
            })

        data = {'slides': slides}
        if not metadata: #kwargs.get('no_metadata'):
            data['metadata'] = self._extract_metadata(presentation)
        return data

    def _extract_all_slides(self, presentation: Presentation, **kwargs) -> Dict[str, Any]:
        slides = []
        for i, slide in enumerate(presentation.slides):
            try:
                slides.append({
                    'slide_number': i + 1,
                    'content': self._extract_slide_content(slide),
                    'notes': self._extract_speaker_notes(slide)
                })
            except Exception as e:
                logger.error(f"Error processing slide {i + 1}: {str(e)}")

        data = {'slides': slides}
        if not kwargs.get('no_metadata'):
            data['metadata'] = self._extract_metadata(presentation)
        return data

    def _extract_slide_content(self, slide) -> str:
        return ' '.join(shape.text for shape in slide.shapes if hasattr(shape, 'text'))

    def _extract_speaker_notes(self, slide) -> str:
        return slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ''

    def _extract_metadata(self, presentation: Presentation) -> Dict[str, Any]:
        core_properties = presentation.core_properties
        return {
            'title': core_properties.title or 'Untitled',
            'author': core_properties.author or 'Unknown',
            'created': str(core_properties.created) if core_properties.created else 'Unknown',
            'modified': str(core_properties.modified) if core_properties.modified else 'Unknown',
            'slide_count': len(presentation.slides)
        }

